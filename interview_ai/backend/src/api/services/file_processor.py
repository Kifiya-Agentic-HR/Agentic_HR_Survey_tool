import PyPDF2
from docx import Document
from pptx import Presentation
from typing import Union
import io

async def process_file(file: Union[bytes, str]) -> str:
    content = await file.read()
    
    if file.filename.endswith('.pdf'):
        return process_pdf(content)
    elif file.filename.endswith('.docx'):
        return process_docx(content)
    elif file.filename.endswith('.pptx'):
        return process_pptx(content)
    else:
        return content.decode('utf-8')

def process_pdf(content: bytes) -> str:
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
    return "\n".join([page.extract_text() for page in pdf_reader.pages])

def process_docx(content: bytes) -> str:
    doc = Document(io.BytesIO(content))
    return "\n".join([para.text for para in doc.paragraphs])

def process_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 100) -> list:
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunks.append(text[start:end])
        start = end - overlap
    return chunks